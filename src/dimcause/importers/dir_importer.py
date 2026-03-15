import hashlib
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from rich.progress import Progress

from dimcause.core.models import Event, EventType, SourceType
from dimcause.storage.markdown_store import MarkdownStore
from dimcause.storage.vector_store import VectorStore


class DirectoryImporter:
    def __init__(
        self,
        base_dir: Path,
        markdown_store: Optional[MarkdownStore] = None,
        vector_store: Optional[VectorStore] = None,
    ):
        self.base_dir = base_dir

        # 依赖注入 / 默认初始化
        self.markdown_store = markdown_store or MarkdownStore()

        # VectorStore
        if vector_store:
            self.vector_store = vector_store
        else:
            try:
                self.vector_store = VectorStore()
            except Exception as e:
                import logging

                logging.getLogger(__name__).warning(
                    f"VectorStore unavailable, semantic search disabled: {e}"
                )
                self.vector_store = None

        # 支持的文件扩展名 (按类型分组)
        self.extensions = {
            # 纯文本 / Markdown
            ".md",
            ".markdown",
            ".txt",
            ".text",
            ".log",
            ".rst",
            ".org",
            # 结构化数据 / 配置
            ".json",
            ".yaml",
            ".yml",
            ".toml",
            ".xml",
            ".csv",
            ".ini",
            ".conf",
            ".env",
            # 代码文件
            ".py",
            ".js",
            ".ts",
            ".tsx",
            ".jsx",
            ".html",
            ".htm",
            ".css",
            ".scss",
            ".java",
            ".go",
            ".rs",
            ".c",
            ".cpp",
            ".h",
            ".hpp",
            ".sh",
            ".zsh",
            ".bash",
            ".sql",
            # Office 文档 (需依赖)
            ".pdf",
            ".docx",
            ".xlsx",
            ".pptx",
        }

    def import_directory(
        self, target_dir: str, recursive: bool = True, pattern: str = "*.md"
    ) -> int:
        """
        导入指定目录下的文件为 Dimcause Events
        """
        target_path = Path(target_dir).resolve()
        if not target_path.exists():
            raise FileNotFoundError(f"Directory not found: {target_dir}")

        if recursive:
            files = [
                p
                for p in target_path.rglob("*")
                if p.suffix.lower() in self.extensions and p.is_file()
            ]
        else:
            files = [
                p
                for p in target_path.glob("*")
                if p.suffix.lower() in self.extensions and p.is_file()
            ]

        files = [f for f in files if ".mal" not in str(f) and ".git" not in str(f)]

        count = 0
        failed_files = []
        with Progress() as progress:
            task = progress.add_task(f"[green]Importing {len(files)} files...", total=len(files))

            for file_path in files:
                try:
                    self._process_file(file_path)
                    count += 1
                except Exception as e:
                    failed_files.append((file_path, str(e)))
                    import logging

                    logging.getLogger(__name__).error(f"Failed to process {file_path}: {e}")
                progress.advance(task)

        # 汇报失败文件
        if failed_files:
            import logging

            logging.getLogger(__name__).warning(
                f"⚠️ {len(failed_files)} files failed to import. "
                f"First 3: {[str(f[0]) for f in failed_files[:3]]}"
            )

        return count

    def _extract_pdf_text(self, file_path: Path) -> str:
        """从 PDF 文件中提取纯文本内容（依赖 pypdf）"""
        try:
            from pypdf import PdfReader
        except ImportError:
            logging.getLogger(__name__).warning(
                f"跳过 PDF 文件 {file_path.name}：未安装 pypdf。"
                "请运行 pip install 'dimcause[full]' 安装 PDF 支持。"
            )
            return ""

        reader = PdfReader(file_path)
        pages_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages_text.append(text)
        return "\n\n".join(pages_text)

    def _extract_docx(self, file_path: Path) -> str:
        """从 Word 文档提取文本 (依赖 python-docx)"""
        try:
            import docx
        except ImportError:
            logging.getLogger(__name__).warning(
                f"跳过 DOCX 文件 {file_path.name}: 未安装 python-docx"
            )
            return ""

        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logging.getLogger(__name__).error(f"DOCX 解析失败 {file_path}: {e}")
            return ""

    def _extract_xlsx(self, file_path: Path) -> str:
        """从 Excel 文档提取文本 (依赖 openpyxl)"""
        try:
            import openpyxl
        except ImportError:
            logging.getLogger(__name__).warning(f"跳过 XLSX 文件 {file_path.name}: 未安装 openpyxl")
            return ""

        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                text_parts.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    # Filter None values and convert to string
                    row_text = "\t".join([str(cell) for cell in row if cell is not None])
                    if row_text:
                        text_parts.append(row_text)
            return "\n".join(text_parts)
        except Exception as e:
            logging.getLogger(__name__).error(f"XLSX 解析失败 {file_path}: {e}")
            return ""

    def _extract_pptx(self, file_path: Path) -> str:
        """从 PowerPoint 文档提取文本 (依赖 python-pptx)"""
        try:
            from pptx import Presentation
        except ImportError:
            logging.getLogger(__name__).warning(
                f"跳过 PPTX 文件 {file_path.name}: 未安装 python-pptx"
            )
            return ""

        try:
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            logging.getLogger(__name__).error(f"PPTX 解析失败 {file_path}: {e}")
            return ""

    def _read_text_file(self, file_path: Path) -> str:
        """尝试多种编码读取文本文件"""
        # 1. 尝试 UTF-8 (最常见)
        try:
            return file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            pass

        # 2. 尝试 GB18030 (涵盖 GBK, GB2312, 常用于中文 Windows)
        try:
            return file_path.read_text(encoding="gb18030")
        except UnicodeDecodeError:
            pass

        # 3. 尝试 Latin-1 (最后手段，保证不崩，但可能是乱码)
        return file_path.read_text(encoding="latin-1", errors="replace")

    def _extract_html(self, file_path: Path) -> str:
        """从 HTML 提取纯文本 (尝试 BeautifulSoup，降级为纯文本)"""
        content = self._read_text_file(file_path)
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(content, "html.parser")
            return soup.get_text(separator="\n")
        except ImportError:
            # 如果没有 bs4，直接返回源码，但在 summary 中可以标记
            return content

    def _process_file(self, file_path: Path):
        suffix = file_path.suffix.lower()

        # 根据类型分发提取逻辑
        if suffix == ".pdf":
            content = self._extract_pdf_text(file_path)
        elif suffix == ".docx":
            content = self._extract_docx(file_path)
        elif suffix == ".xlsx":
            content = self._extract_xlsx(file_path)
        elif suffix == ".pptx":
            content = self._extract_pptx(file_path)
        elif suffix in {".html", ".htm"}:
            content = self._extract_html(file_path)
        else:
            # 默认作为纯文本读取
            content = self._read_text_file(file_path)

        if not content.strip():
            return

        # Generate stable ID based on file path (so re-import updates instead of duplicates)
        file_hash = hashlib.md5(str(file_path).encode()).hexdigest()[:12]
        event_id = f"doc_{file_hash}"

        # Try to parse existing frontmatter if any
        # If not, use file stats
        stats = file_path.stat()
        datetime.fromtimestamp(stats.st_ctime)
        modified_at = datetime.fromtimestamp(stats.st_mtime)

        # Use python-frontmatter if available to get title/metadata
        metadata = {}
        summary = file_path.stem

        try:
            import frontmatter

            post = frontmatter.loads(content)
            if post.metadata:
                metadata.update(post.metadata)
                content = post.content  # Use content without frontmatter
                # Try to find title in metadata
                if "title" in metadata:
                    summary = metadata["title"]
        except ImportError:
            pass

        # Fallback: Try to find title from first line if still default
        if summary == file_path.stem:
            lines = [line.strip() for line in content.splitlines() if line.strip()]
            if lines:
                candidate = lines[0].lstrip("#").strip()
                if len(candidate) < 100:
                    summary = candidate

        # Create Event
        event = Event(
            id=event_id,
            type=EventType.RESOURCE,  # Treat external docs as Resources
            source=SourceType.FILE,
            timestamp=modified_at,
            summary=summary,
            content=content,
            metadata={
                **metadata,
                "original_path": str(file_path),
                "imported_at": datetime.now().isoformat(),
            },
        )

        # Save
        self.markdown_store.save(event)  # This will copy it to .mal/events

        if self.vector_store:
            # Add to vector store for semantic search
            self.vector_store.add(event)


def run_dir_import(path: str, recursive: bool = True):
    importer = DirectoryImporter(Path("."))
    count = importer.import_directory(path, recursive)
    print(f"✅ Successfully imported {count} documents from {path}")
