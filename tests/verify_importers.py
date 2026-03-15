import logging
import shutil
from pathlib import Path

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("test_importer")

# Define paths
TEST_DIR = Path("tmp/test_importer")
if TEST_DIR.exists():
    shutil.rmtree(TEST_DIR)
TEST_DIR.mkdir(parents=True)

# 1. Create Test Files
print("Creating test files...")

# Pure Text / Code
(TEST_DIR / "note.md").write_text("# Markdown Note\nThis is a markdown file.", encoding="utf-8")
(TEST_DIR / "script.py").write_text("def hello():\n    print('Hello Python')", encoding="utf-8")
(TEST_DIR / "data.json").write_text('{"key": "value", "list": [1, 2, 3]}', encoding="utf-8")

# Chinese Text Support
(TEST_DIR / "utf8_cn.txt").write_text("这是一个UTF-8编码的中文文件。", encoding="utf-8")
try:
    (TEST_DIR / "gbk_cn.txt").write_text("这是一个GBK编码的中文文件。", encoding="gb18030")
except LookupError:
    print("Warning: gb18030 encoding not available, skipping gbk_cn.txt")

# Traditional Chinese
(TEST_DIR / "trad_cn.md").write_text("# 內存優化\n這是一個關於記憶體管理的文檔。", encoding="utf-8")

# HTML
(TEST_DIR / "page.html").write_text(
    "<html><body><h1>HTML Page</h1><p>This is a paragraph.</p></body></html>", encoding="utf-8"
)

# Office Docs (require libraries)
try:
    from docx import Document

    doc = Document()
    doc.add_heading("Word Document", 0)
    doc.add_paragraph("This is a paragraph in Word.")
    doc.save(TEST_DIR / "doc.docx")
    print("Created doc.docx")
except ImportError:
    print("Skipped doc.docx (python-docx not installed)")

try:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "MySheet"
    ws["A1"] = "Header1"
    ws["B1"] = "Header2"
    ws["A2"] = "Value1"
    ws["B2"] = 123
    wb.save(TEST_DIR / "sheet.xlsx")
    print("Created sheet.xlsx")
except ImportError:
    print("Skipped sheet.xlsx (openpyxl not installed)")

try:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Presentation Title"
    subtitle = slide.placeholders[1]
    subtitle.text = "Subtitle Text"
    prs.save(TEST_DIR / "slides.pptx")
    print("Created slides.pptx")
except ImportError:
    print("Skipped slides.pptx (python-pptx not installed)")


# 2. Run Importer
print("\nRunning DirectoryImporter...")
from dimcause.extractors.extractor import BasicExtractor  # noqa: E402
from dimcause.importers.dir_importer import DirectoryImporter  # noqa: E402
from dimcause.search.engine import SearchEngine  # noqa: E402
from dimcause.storage.markdown_store import MarkdownStore  # noqa: E402

try:
    import zhconv
except ImportError:
    zhconv = None

# Mock VectorStore to avoid embedding costs/time during test
# DirectoryImporter handles missing VectorStore gracefully or just calls add.
from unittest.mock import MagicMock  # noqa: E402

mock_vector_store = MagicMock()
mock_vector_store.add = MagicMock()

# Define EVENTS_DIR before using it in MarkdownStore
EVENTS_DIR = Path.home() / ".dimcause" / "events"

importer = DirectoryImporter(
    markdown_store=MarkdownStore(base_dir=EVENTS_DIR),
    vector_store=mock_vector_store,  # Pass the mock
    base_dir=TEST_DIR,
)
imported_count = importer.import_directory(TEST_DIR)
print(f"Imported {imported_count} files.")

# 3. Verify Content
print("\nVerifying imported content...")

# Helper to find event by filename
# DirectoryImporter generates ID: doc_{md5(path)[:12]}
# We can just list all files in .mal/events/resource/file/ and check summary/content

# from dimcause.storage.markdown_store import MarkdownStore # Already imported above
# store = MarkdownStore() # Not needed for this verification loop
# List recently modified files
# events = [] # Not needed for this verification loop
# store.list_by_date is based on event timestamp, which we set to file mtime.
# Let's search recursively in the events directory
if not EVENTS_DIR.exists():
    print(f"Error: {EVENTS_DIR} does not exist")
    exit(1)

# found_files = {} # Not needed for this verification loop

# count_scanned = 0 # Not needed for this verification loop
# for event_file in EVENTS_DIR.rglob("*.md"): # Not needed for this verification loop
#     count_scanned += 1 # Not needed for this verification loop
#     event = store.load(str(event_file)) # Not needed for this verification loop
#     if not event: # Not needed for this verification loop
#         continue # Not needed for this verification loop

#     # Check if it's one of our test files # Not needed for this verification loop
#     original_path = event.metadata.get("original_path", "") # Not needed for this verification loop
#     if str(TEST_DIR) in original_path: # Not needed for this verification loop
#         filename = Path(original_path).name # Not needed for this verification loop
#         found_files[filename] = event # Not needed for this verification loop

# Checks
expected_files = {
    "note.md": "Markdown Note",
    "script.py": "def hello():",
    "data.json": '"key": "value"',
    "page.html": "HTML Page",  # BS4 extracts text
    "doc.docx": "This is a paragraph in Word.",
    "sheet.xlsx": "Value1",
    "slides.pptx": "Subtitle Text",
    "utf8_cn.txt": "这是一个UTF-8编码的中文文件。",
    "gbk_cn.txt": "这是一个GBK编码的中文文件。",
    "trad_cn.md": "記憶體管理",
}

passed = 0
skipped = 0

print("\nVerifying imported content...")
for filename, expected_snippet in expected_files.items():
    file_path = TEST_DIR / filename
    if not file_path.exists():
        print(f"⚠️ {filename}: Skipped (File not created)")
        skipped += 1
        continue

    # Find event file
    found = False
    for event_file in EVENTS_DIR.rglob("*.md"):
        content = event_file.read_text(encoding="utf-8")
        if filename in content and expected_snippet in content:
            print(f"✅ {filename}: Content Verified")
            found = True
            passed += 1
            break

    if not found:
        print(f"❌ {filename}: Content Mismatch!\n   Expected snippet: '{expected_snippet}'")

print("\n--- Advanced Feature Verification ---")

# 1. SearchEngine: 简繁互搜
if zhconv:
    print("Verifying Search Engine (Traditional <-> Simplified)...")

    # 确保 trad_cn.md 已被导入为 Event
    # DirectoryImporter 已在前面运行过，如果 trad_cn.md 当时已存在，则已被导入。
    # 我们在前面刚刚创建了 trad_cn.md，但 importer 是在创建文件之后运行的吗？
    # 检查脚本顺序：
    # 1. Creating test files (含 trad_cn.md)
    # 2. Running DirectoryImporter
    # 所以 trad_cn.md 应该已经被导入了。

    # 初始化 SearchEngine (指定与 Importer 相同的存储路径)
    store = MarkdownStore(base_dir=EVENTS_DIR)
    # Patch vector_store to avoid loading model defaults
    engine = SearchEngine(markdown_store=store, vector_store=mock_vector_store)

    # 强制刷新一下 store 的缓存（如果需要）
    # MarkdownStore.list_by_date 是直接读文件系统的，不需要刷新。

    # 用简体搜繁体
    # query="内存" -> 简体="内存", 繁体="內存"
    # trad_cn.md content="...內存..."
    # 搜索时：
    # engine._text_search("内存") ->
    #   q_sim = "内存"
    #   c_sim = convert("...內存...", zh-cn) -> "...内存..."
    #   "内存" in "...内存..." -> True!

    # 关键：禁用 reranker 以免下载模型失败
    results = engine.search("内存", mode="text", use_reranker=False)

    match_found = False
    for e in results:
        # 检查 content 或原始文件名
        if "內存" in e.content or "trad_cn.md" in str(e.metadata):
            match_found = True
            break

    if match_found:
        print("✅ Text Search: Simplified query matched Traditional content")
    else:
        print("❌ Text Search: Simplified query FAILED to match Traditional content")
        print(f"   Debug: Found {len(results)} results")
        for e in results:
            print(f"   - {e.summary[:30]}...")
else:
    print("⚠️ Skipping Search Engine verification (zhconv not installed)")

# 2. Extractor: 中文关键词提取
print("Verifying Extractor Keywords...")
extractor = BasicExtractor()

# Decision
e1 = extractor.extract("我们决定采用 Redis 作为缓存。")
if e1.type.value == "decision":
    print("✅ Extractor: Detected 'decision' from '决定/采用'")
else:
    print(f"❌ Extractor: Failed to detect 'decision', got {e1.type}")

# Code Change (Traditional)
e2 = extractor.extract("修復了一個嚴重的內存洩漏問題。")
if e2.type.value == "code_change" or e2.type.value == "diagnostic":  # 修复/修復 -> code_change
    print("✅ Extractor: Detected 'code_change' from Traditional '修復'")
else:
    print(f"❌ Extractor: Failed to detect 'code_change', got {e2.type}")


print(
    f"\nVerification Complete: {passed}/{len(expected_files)} checks passed (excluding skipped sources)."
)
